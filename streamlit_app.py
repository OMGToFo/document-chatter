import streamlit as st
from openai import OpenAI
import pandas as pd
from docx import Document
import PyPDF2
import io
from pptx import Presentation

# Show title and description.
st.title("📄 Chat with Documents")
st.write(
    "Upload one or more documenta below and ask a questions – GPT will answer! "
    "To use this app, you need to provide an OpenAI API key, which you can get [here](https://platform.openai.com/account/api-keys). "
)

# Ask user for their OpenAI API key via `st.text_input`.
# Alternatively, you can store the API key in `./.streamlit/secrets.toml` and access it
# via `st.secrets`, see https://docs.streamlit.io/develop/concepts/connections/secrets-management
openai_api_key = st.text_input("OpenAI API Key", type="password")
if not openai_api_key:
    st.info("Please add your OpenAI API key to continue.", icon="🗝️")
    st.stop()
else:

    # Create an OpenAI client.
    client = OpenAI(api_key=openai_api_key)

    # Let the user upload files via `st.file_uploader`.
uploaded_files = st.file_uploader(
    "Upload documents (.txt, .md, .doc, .docx, .xls, .xlsx, .pdf, .pptx)",
    type=("txt", "md", "doc", "docx", "xls", "xlsx", "pdf", "pptx"),accept_multiple_files=True
)

question = st.text_area(
    "Now ask a question:",
    placeholder="Can you give me a short summary?",
    disabled=not uploaded_files
)

if uploaded_files and question:

    combined_document = ""
    for uploaded_file in uploaded_files:
        if uploaded_file.name.endswith(".txt") or uploaded_file.name.endswith(".md"):
            combined_document += uploaded_file.read().decode() + "\n"

        elif uploaded_file.name.endswith(".doc") or uploaded_file.name.endswith(".docx"):
            doc = Document(uploaded_file)
            combined_document += "\n".join([para.text for para in doc.paragraphs]) + "\n"

        elif uploaded_file.name.endswith(".xls") or uploaded_file.name.endswith(".xlsx"):
            df = pd.read_excel(uploaded_file, engine='openpyxl')
            combined_document += df.to_string() + "\n"

        elif uploaded_file.name.endswith(".pdf"):
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                combined_document += page.extract_text() + "\n"

        elif uploaded_file.name.endswith(".pptx"):
            presentation = Presentation(uploaded_file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        combined_document += shape.text + "\n"

    messages = [
        {
            "role": "user",
            "content": f"Here are the documents: {combined_document} \n\n---\n\n {question}",
        }
    ]

    try:
        stream = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=messages,
            stream=True,
        )
    
        st.write_stream(stream)

    except Exception as e:
        st.error(f"An error occurred: {str(e)}")

